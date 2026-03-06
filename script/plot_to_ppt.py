#!/usr/bin/env python3
"""
plot_to_ppt.py
===============
Generate a PowerPoint report containing the raw statistics and graphs produced
by `plot_unified_geomean.py`.

Usage:
  python3 plot_to_ppt.py

Requirements:
  pip install python-pptx

This script re-uses the data-gathering, processing, and plotting utilities
from `plot_unified_geomean.py`.  It first produces the usual PNG graphs (if
they are not already present) and then assembles them—together with tables of
raw geomean statistics—into a single PowerPoint file located in the `graphs/`
directory.
"""

# TITLE = "CVP_server"
TITLE = ""

import os
import sys
from importlib import import_module
from typing import List, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
except ImportError as exc:  # pragma: no cover
    sys.exit("python-pptx is required.  Install it via `pip install python-pptx`.  Aborting.")

# ---------------------------------------------------------------------------
# Import the companion module *without* executing its `main()`
# ---------------------------------------------------------------------------
# The companion script lives in the same directory as this file.  We add that
# directory to ``sys.path`` to make the import straightforward regardless of
# whether ``scripts`` is a proper Python package.

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
if SCRIPT_DIR not in sys.path:
    sys.path.insert(0, SCRIPT_DIR)

try:
    import plot_unified_geomean as pu  # type: ignore
except ModuleNotFoundError as exc:
    sys.exit("Unable to import 'plot_unified_geomean.py'.  Make sure it is located in the same "
             "directory as this script. Aborting.")

# ---------------------------------------------------------------------------
# Helper functions for building the deck
# ---------------------------------------------------------------------------


def _add_title_slide(prs: Presentation, title: str) -> None:
    """Insert a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    slide.shapes.title.text = title


def _add_image_slide(prs: Presentation, title: str, image_path: str) -> None:
    """Insert a slide with a single image."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    # Title textbox
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    txbox.text_frame.text = title
    # The picture itself
    slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.0), width=Inches(9))


def _add_table_slide(
    prs: Presentation,
    title: str,
    headers: List[str],
    rows: List[Tuple[str, str]],
) -> None:
    """Insert a slide that contains a simple table."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = txbox.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(24)

    # Table
    rows_n = len(rows) + 1  # +1 for header
    cols_n = len(headers)
    table_shape = slide.shapes.add_table(
        rows_n,
        cols_n,
        Inches(0.5),
        Inches(1.0),
        Inches(9),
        Inches(0.5 + 0.3 * rows_n),
    )
    table = table_shape.table

    # Header row
    for col_idx, header in enumerate(headers):
        table.cell(0, col_idx).text = header

    # Data rows
    for row_idx, data_row in enumerate(rows, start=1):
        for col_idx, value in enumerate(data_row):
            table.cell(row_idx, col_idx).text = str(value)

    # Center align all cells
    for row in table.rows:
        for cell in row.cells:
            cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


# ---------------------------------------------------------------------------
# Utility helpers
# ---------------------------------------------------------------------------


def _speedups_to_rows(speed_dict, pfetchers):
    """Convert geomean dictionary to a list of (prefetcher, value) tuples."""
    return [
        (pf, f"{speed_dict[pf].get('geomean', 0.0):.5f}") for pf in pfetchers
    ]


# ---------------------------------------------------------------------------
# Core slide-generation routine for a single benchmark suite
# ---------------------------------------------------------------------------


_SPEC_EXCLUDE = {}


def _add_suite_to_presentation(
    prs: Presentation,
    benchmarks,
    suite_label: str,
    base_plot_name: str,
    benchmark_weights=None,
    benchmark_shortcode=None,
    spec_filter: bool = True,
) -> None:  # noqa: C901 – fine
    """Generate plots & slides for *benchmarks* (list) and append to *prs*.

    Parameters
    ----------
    benchmark_weights:
        If provided, overrides ``pu.BENCHMARK_WEIGHTS`` for this suite.
    benchmark_shortcode:
        If provided, overrides ``pu.BENCHMARK_SHORTCODE`` for this suite.
    spec_filter:
        When *True* (default) the SPEC-specific problematic benchmarks
        (milc433, lbm470, gromacs435, exchange2648) are removed from the list.
        Set to *False* for non-SPEC suites such as GAP and CVP_SRV.
    """

    # Backup original global settings so we can restore afterwards
    original_benchmarks = pu.BENCHMARKS.copy()
    original_plot_name = pu.PLOT_NAME
    original_weights = pu.BENCHMARK_WEIGHTS
    original_shortcode = pu.BENCHMARK_SHORTCODE

    # Update globals in plot_unified_geomean
    filtered = [bm for bm in benchmarks if bm not in _SPEC_EXCLUDE] if spec_filter else list(benchmarks)
    pu.BENCHMARKS = filtered
    pu.PLOT_NAME = f"{suite_label}_{base_plot_name}"
    if benchmark_weights is not None:
        pu.BENCHMARK_WEIGHTS = benchmark_weights
    if benchmark_shortcode is not None:
        pu.BENCHMARK_SHORTCODE = benchmark_shortcode

    # ----------------------------
    # Data collection / processing
    # ----------------------------
    print(f"[PPT] Gathering raw data for {suite_label}…")
    ipc_data = pu.gather_data(pu.parse_ipc_from_file, "ipc")
    dram_data = pu.gather_data(pu.parse_dram_from_file, "dram")
    cov_data = pu.gather_data(pu.parse_cov_from_file, "coverage")
    acc_data = pu.gather_data(pu.parse_acc_from_file, "accuracy")

    ipc_speedups, ipc_plot_pref = pu.compute_geomean_speedups(ipc_data, "ipc", pu.BASELINE)
    dram_speedups, dram_plot_pref = pu.compute_geomean_speedups(dram_data, "dram", pu.BASELINE)
    cov_speedups, cov_plot_pref = pu.compute_geomean_speedups(cov_data, "coverage", pu.BASELINE)
    acc_speedups, acc_plot_pref = pu.compute_geomean_speedups(acc_data, "accuracy", pu.BASELINE)

    # ----------------------------
    # Ensure the graphs exist
    # ----------------------------
    print(f"[PPT] Generating graphs for {suite_label}…")

    def _ensure_plot(metric, ylabel, speedups, plot_pref):
        pu.create_plot(
            speedups,
            plot_pref,
            metric,
            ylabel,
            f"{pu.PLOT_NAME}_{metric}.{pu.OUTPUT}",
            include_geomean=pu.INCLUDE_GEOMEAN,
            include_baseline=(metric in {"ipc", "dram"}),
            ylim_bottom=0.9 if metric in {"ipc", "dram"} else 0.0,
            ylim_top=None if metric in {"ipc", "dram"} else 1.0,
            only_geomean_bar=pu.ONLY_GEOMEAN_BAR,
            legend_position="top",
        )

    _ensure_plot("ipc", "IPC Speedup", ipc_speedups, ipc_plot_pref)
    _ensure_plot("dram", "Normalized DRAM Traffic", dram_speedups, dram_plot_pref)
    _ensure_plot("coverage", "Coverage", cov_speedups, cov_plot_pref)
    _ensure_plot("accuracy", "Accuracy", acc_speedups, acc_plot_pref)

    # ----------------------------
    # Append slides
    # ----------------------------
    _add_title_slide(prs, f"{suite_label} – {TITLE}")

    _add_table_slide(
        prs,
        "IPC – Geometric Mean Speedup",
        ["Prefetcher", "Geo-Speedup"],
        _speedups_to_rows(ipc_speedups, ipc_plot_pref),
    )
    _add_table_slide(
        prs,
        "DRAM – Normalised Traffic",
        ["Prefetcher", "Normalised"],
        _speedups_to_rows(dram_speedups, dram_plot_pref),
    )
    _add_table_slide(
        prs,
        "Coverage",
        ["Prefetcher", "Coverage"],
        _speedups_to_rows(cov_speedups, cov_plot_pref),
    )
    _add_table_slide(
        prs,
        "Accuracy",
        ["Prefetcher", "Accuracy"],
        _speedups_to_rows(acc_speedups, acc_plot_pref),
    )

    img_titles = ["IPC Speedup", "DRAM Traffic", "Coverage", "Accuracy"]
    img_files = [
        f"{pu.PLOT_NAME}_ipc.{pu.OUTPUT}",
        f"{pu.PLOT_NAME}_dram.{pu.OUTPUT}",
        f"{pu.PLOT_NAME}_coverage.{pu.OUTPUT}",
        f"{pu.PLOT_NAME}_accuracy.{pu.OUTPUT}",
    ]

    for title, fname in zip(img_titles, img_files):
        path = os.path.join(pu.GRAPH_DIR, fname)
        if os.path.exists(path):
            _add_image_slide(prs, title, path)
        else:
            print(f"[PPT] Warning: graph not found – {path}")

    # ----------------------------
    # Restore globals
    # ----------------------------
    pu.BENCHMARKS = original_benchmarks
    pu.PLOT_NAME = original_plot_name
    pu.BENCHMARK_WEIGHTS = original_weights
    pu.BENCHMARK_SHORTCODE = original_shortcode


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------


def main() -> None:  # noqa: C901 – length OK for script
    base_plot_name = pu.PLOT_NAME  # e.g., misc
    benchmark_type = pu.BENCHMARK_TYPE  # driven by plot_unified_geomean.py

    prs = Presentation()

    # Dispatch to the single suite selected by BENCHMARK_TYPE
    if benchmark_type == 'SPEC2006':
        _add_suite_to_presentation(
            prs, pu.spec2006_ones, "SPEC2006", base_plot_name, spec_filter=True,
        )
    elif benchmark_type == 'SPEC2017':
        _add_suite_to_presentation(
            prs, pu.spec2017_ones, "SPEC2017", base_plot_name, spec_filter=True,
        )
    elif benchmark_type == 'SPEC_ALL':
        _add_suite_to_presentation(
            prs, list(pu.SPEC2017_shortcode.keys()), "SPEC_ALL", base_plot_name, spec_filter=True,
        )
    elif benchmark_type == 'GAP':
        _add_suite_to_presentation(
            prs, pu.gap_ones, "GAP", base_plot_name,
            benchmark_weights=pu.GAP_SHORTCODE_WEIGHTS,
            benchmark_shortcode=pu.GAP_shortcode,
            spec_filter=False,
        )
    elif benchmark_type == 'CVP_SRV':
        _add_suite_to_presentation(
            prs, pu.cvp_srv_ones, "CVP_SRV", base_plot_name,
            benchmark_weights=pu.CVP_SRV_SHORTCODE_WEIGHTS,
            benchmark_shortcode=pu.CVP_SRV_shortcode,
            spec_filter=False,
        )
    elif benchmark_type == 'NEW_GAP':
        _add_suite_to_presentation(
            prs, pu.new_gap_ones, "NEW_GAP", base_plot_name,
            benchmark_weights=pu.New_GAP_SHORTCODE_WEIGHTS,
            benchmark_shortcode=pu.New_GAP_shortcode,
            spec_filter=False,
        )
    elif benchmark_type == 'GOOGLE':
        _add_suite_to_presentation(
            prs, pu.google_traces_ones, "GOOGLE", base_plot_name,
            benchmark_weights=pu.Google_Traces_SHORTCODE_WEIGHTS,
            benchmark_shortcode=pu.Google_Traces_shortcode,
            spec_filter=False,
        )
    else:
        sys.exit(f"Unknown BENCHMARK_TYPE '{benchmark_type}' in plot_unified_geomean.py. Aborting.")

    # Save the deck
    os.makedirs(pu.GRAPH_DIR, exist_ok=True)
    ppt_path = os.path.join(pu.GRAPH_DIR, f"{base_plot_name}_{benchmark_type}_summary.pptx")
    prs.save(ppt_path)
    print(f"[PPT] Presentation saved to: {ppt_path}")


if __name__ == "__main__":  # pragma: no cover
    main()
    # Convert PPTX to PDF using libreoffice, if available
    try:
        import shutil
        ppt_path = os.path.join(pu.GRAPH_DIR, f"{pu.PLOT_NAME}_{pu.BENCHMARK_TYPE}_summary.pptx")
        pdf_path = os.path.splitext(ppt_path)[0] + ".pdf"

        # Check if libreoffice is available
        soffice_path = shutil.which("libreoffice") or shutil.which("soffice")
        if soffice_path is not None:
            import subprocess
            cmd = [
                soffice_path, '--headless', '--convert-to', 'pdf',
                '--outdir', pu.GRAPH_DIR, ppt_path
            ]
            subprocess.run(cmd, check=True)
            if os.path.exists(pdf_path):
                print(f"[PPT] PDF also saved to: {pdf_path}")
            else:
                print(f"[PPT] Tried converting to PDF but output file not found: {pdf_path}")
        else:
            print("[PPT] LibreOffice ('libreoffice' or 'soffice') not found in PATH; skipping PDF export.")
    except Exception as e:
        print(f"[PPT] Error during PDF export: {e}")
